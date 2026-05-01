import streamlit as st
import io
import os
import subprocess
import tempfile
import zipfile
from pathlib import Path

# ─── Library availability ───────────────────────────────────────────────────────
def _has(pkg):
    try:
        __import__(pkg)
        return True
    except ImportError:
        return False

HAS_PYMUPDF   = _has("fitz")
HAS_PDF2DOCX  = _has("pdf2docx")
HAS_PPTX      = _has("pptx")
HAS_DOCX      = _has("docx")
HAS_PIL       = _has("PIL")
HAS_REPORTLAB = _has("reportlab")
HAS_PANDAS    = _has("pandas") and _has("openpyxl")


# ─── Format config ──────────────────────────────────────────────────────────────
CONVERSIONS = {
    "PDF":  ["DOCX", "TXT", "PNG", "JPG", "PPTX"],
    "DOCX": ["PDF", "TXT"],
    "PPTX": ["PDF", "PNG", "JPG"],
    "PNG":  ["JPG", "PDF", "WEBP", "BMP", "DOCX"],
    "JPG":  ["PNG", "PDF", "WEBP", "BMP", "DOCX"],
    "WEBP": ["PNG", "JPG", "PDF"],
    "BMP":  ["PNG", "JPG", "PDF", "WEBP"],
    "TXT":  ["PDF", "DOCX"],
    "CSV":  ["XLSX", "JSON"],
    "XLSX": ["CSV", "JSON"],
    "JSON": ["CSV", "XLSX"],
}

ACCEPT = {
    "PDF":  ["pdf"],
    "DOCX": ["docx"],
    "PPTX": ["pptx"],
    "PNG":  ["png"],
    "JPG":  ["jpg", "jpeg"],
    "WEBP": ["webp"],
    "BMP":  ["bmp"],
    "TXT":  ["txt"],
    "CSV":  ["csv"],
    "XLSX": ["xlsx", "xls"],
    "JSON": ["json"],
}

ICONS = {
    "PDF": "📄", "DOCX": "📝", "PPTX": "📊",
    "PNG": "🖼️", "JPG": "🖼️", "WEBP": "🖼️", "BMP": "🖼️",
    "TXT": "📃", "CSV": "📋", "XLSX": "📈", "JSON": "🗂️",
}

MIME = {
    "PDF":  "application/pdf",
    "DOCX": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "PPTX": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "PNG":  "image/png",
    "JPG":  "image/jpeg",
    "WEBP": "image/webp",
    "BMP":  "image/bmp",
    "TXT":  "text/plain",
    "CSV":  "text/csv",
    "XLSX": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "JSON": "application/json",
    "ZIP":  "application/zip",
}

REQUIRES = {
    ("PDF",  "DOCX"): ["PyMuPDF", "pdf2docx"],
    ("PDF",  "TXT"):  ["PyMuPDF"],
    ("PDF",  "PNG"):  ["PyMuPDF"],
    ("PDF",  "JPG"):  ["PyMuPDF"],
    ("PDF",  "PPTX"): ["PyMuPDF", "python-pptx"],
    ("DOCX", "PDF"):  ["LibreOffice"],
    ("DOCX", "TXT"):  ["python-docx"],
    ("PPTX", "PDF"):  ["LibreOffice"],
    ("PPTX", "PNG"):  ["LibreOffice", "PyMuPDF"],
    ("PPTX", "JPG"):  ["LibreOffice", "PyMuPDF"],
    ("PNG",  "JPG"):  ["Pillow"],
    ("PNG",  "PDF"):  ["Pillow"],
    ("PNG",  "WEBP"): ["Pillow"],
    ("PNG",  "BMP"):  ["Pillow"],
    ("PNG",  "DOCX"): ["Pillow", "python-docx"],
    ("JPG",  "PNG"):  ["Pillow"],
    ("JPG",  "PDF"):  ["Pillow"],
    ("JPG",  "WEBP"): ["Pillow"],
    ("JPG",  "BMP"):  ["Pillow"],
    ("JPG",  "DOCX"): ["Pillow", "python-docx"],
    ("WEBP", "PNG"):  ["Pillow"],
    ("WEBP", "JPG"):  ["Pillow"],
    ("WEBP", "PDF"):  ["Pillow"],
    ("BMP",  "PNG"):  ["Pillow"],
    ("BMP",  "JPG"):  ["Pillow"],
    ("BMP",  "PDF"):  ["Pillow"],
    ("BMP",  "WEBP"): ["Pillow"],
    ("TXT",  "PDF"):  ["reportlab"],
    ("TXT",  "DOCX"): ["python-docx"],
    ("CSV",  "XLSX"): ["pandas", "openpyxl"],
    ("CSV",  "JSON"): ["pandas"],
    ("XLSX", "CSV"):  ["pandas", "openpyxl"],
    ("XLSX", "JSON"): ["pandas", "openpyxl"],
    ("JSON", "CSV"):  ["pandas"],
    ("JSON", "XLSX"): ["pandas", "openpyxl"],
}


# ─── LibreOffice helpers ────────────────────────────────────────────────────────
def _find_libreoffice():
    candidates = [
        "libreoffice", "soffice",
        "/usr/bin/libreoffice", "/usr/bin/soffice",
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ]
    for c in candidates:
        try:
            subprocess.run([c, "--version"], capture_output=True, timeout=5)
            return c
        except Exception:
            continue
    return None


def _lo_convert(data, in_ext, out_ext):
    lo = _find_libreoffice()
    if not lo:
        raise RuntimeError(
            "LibreOffice is required for this conversion but was not found.\n\n"
            "Install it:\n"
            "  macOS  →  brew install --cask libreoffice\n"
            "  Ubuntu →  sudo apt install libreoffice\n"
            "  Windows → https://www.libreoffice.org/download"
        )
    with tempfile.TemporaryDirectory() as tmp:
        src = Path(tmp) / f"input.{in_ext}"
        src.write_bytes(data)
        r = subprocess.run(
            [lo, "--headless", "--convert-to", out_ext, "--outdir", tmp, str(src)],
            capture_output=True, timeout=180
        )
        out = Path(tmp) / f"input.{out_ext}"
        if out.exists():
            return out.read_bytes()
        raise RuntimeError(
            f"LibreOffice conversion failed.\n"
            f"stdout: {r.stdout.decode(errors='replace')}\n"
            f"stderr: {r.stderr.decode(errors='replace')}"
        )


# ─── PDF conversions ────────────────────────────────────────────────────────────
def pdf_to_docx(data):
    from pdf2docx import Converter
    with tempfile.TemporaryDirectory() as tmp:
        src = Path(tmp) / "input.pdf"
        dst = Path(tmp) / "output.docx"
        src.write_bytes(data)
        cv = Converter(str(src))
        cv.convert(str(dst))
        cv.close()
        return dst.read_bytes(), "converted.docx", MIME["DOCX"]


def pdf_to_txt(data):
    import fitz
    doc = fitz.open(stream=data, filetype="pdf")
    pages = []
    for i, page in enumerate(doc):
        pages.append(f"--- Page {i+1} ---\n{page.get_text()}")
    doc.close()
    return "\n\n".join(pages).encode("utf-8"), "converted.txt", MIME["TXT"]


def pdf_to_image(data, fmt):
    import fitz
    fmt_lower = fmt.lower()
    pil_fmt = "jpeg" if fmt == "JPG" else fmt_lower
    doc = fitz.open(stream=data, filetype="pdf")
    mat = fitz.Matrix(2.0, 2.0)  # 2x = 144 DPI — crisp output

    if len(doc) == 1:
        pix = doc[0].get_pixmap(matrix=mat)
        img_bytes = pix.tobytes(output=pil_fmt)
        doc.close()
        return img_bytes, f"converted.{fmt_lower}", MIME[fmt]

    # Multi-page → ZIP of images
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for i, page in enumerate(doc):
            pix = page.get_pixmap(matrix=mat)
            zf.writestr(f"page_{i+1:03d}.{fmt_lower}", pix.tobytes(output=pil_fmt))
    doc.close()
    buf.seek(0)
    return buf.read(), "pages.zip", MIME["ZIP"]


def pdf_to_pptx(data):
    import fitz
    from pptx import Presentation
    from pptx.util import Emu
    doc = fitz.open(stream=data, filetype="pdf")
    prs = Presentation()
    r = doc[0].rect
    # PDF points → EMU (1 pt = 12700 EMU)
    prs.slide_width  = Emu(int(r.width  * 12700))
    prs.slide_height = Emu(int(r.height * 12700))
    blank_layout = prs.slide_layouts[6]
    mat = fitz.Matrix(2.0, 2.0)
    for page in doc:
        pix = page.get_pixmap(matrix=mat)
        img_buf = io.BytesIO(pix.tobytes(output="png"))
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            img_buf, 0, 0,
            width=prs.slide_width, height=prs.slide_height
        )
    doc.close()
    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out.read(), "converted.pptx", MIME["PPTX"]


# ─── DOCX conversions ───────────────────────────────────────────────────────────
def docx_to_pdf(data):
    return _lo_convert(data, "docx", "pdf"), "converted.pdf", MIME["PDF"]


def docx_to_txt(data):
    from docx import Document
    from docx.oxml.ns import qn
    doc = Document(io.BytesIO(data))
    lines = []
    for para in doc.paragraphs:
        lines.append(para.text)
    # Also grab table text
    for table in doc.tables:
        for row in table.rows:
            lines.append("\t".join(cell.text for cell in row.cells))
    return "\n".join(lines).encode("utf-8"), "converted.txt", MIME["TXT"]


# ─── PPTX conversions ───────────────────────────────────────────────────────────
def pptx_to_pdf(data):
    return _lo_convert(data, "pptx", "pdf"), "converted.pdf", MIME["PDF"]


def pptx_to_image(data, fmt):
    # PPTX → PDF → images
    pdf_data = _lo_convert(data, "pptx", "pdf")
    return pdf_to_image(pdf_data, fmt)


# ─── Image conversions ──────────────────────────────────────────────────────────
def img_to_img(data, to_fmt):
    from PIL import Image
    img = Image.open(io.BytesIO(data))
    # RGBA/P not supported in JPG/BMP
    if to_fmt in ("JPG", "BMP") and img.mode in ("RGBA", "P", "LA"):
        img = img.convert("RGB")
    buf = io.BytesIO()
    pil_fmt = "JPEG" if to_fmt == "JPG" else to_fmt
    img.save(buf, format=pil_fmt)
    buf.seek(0)
    ext = "jpg" if to_fmt == "JPG" else to_fmt.lower()
    return buf.read(), f"converted.{ext}", MIME[to_fmt]


def img_to_pdf(data):
    from PIL import Image
    img = Image.open(io.BytesIO(data))
    if img.mode in ("RGBA", "P", "LA"):
        img = img.convert("RGB")
    buf = io.BytesIO()
    img.save(buf, format="PDF")
    buf.seek(0)
    return buf.read(), "converted.pdf", MIME["PDF"]


def img_to_docx(data):
    from docx import Document
    from docx.shared import Inches
    from PIL import Image
    img = Image.open(io.BytesIO(data))
    # Convert to PNG in-memory if needed (DOCX supports PNG/JPEG reliably)
    png_buf = io.BytesIO()
    if img.mode in ("RGBA", "P", "LA"):
        img = img.convert("RGBA")
    elif img.mode != "RGB":
        img = img.convert("RGB")
    img.save(png_buf, format="PNG")
    png_buf.seek(0)
    doc = Document()
    # Fit to 6-inch width max
    doc.add_picture(png_buf, width=Inches(6))
    out = io.BytesIO()
    doc.save(out)
    out.seek(0)
    return out.read(), "converted.docx", MIME["DOCX"]


# ─── TXT conversions ────────────────────────────────────────────────────────────
def txt_to_pdf(data):
    from reportlab.lib.pagesizes import A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.enums import TA_LEFT
    text = data.decode("utf-8", errors="replace")
    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=A4,
        leftMargin=72, rightMargin=72, topMargin=72, bottomMargin=72
    )
    styles = getSampleStyleSheet()
    style = styles["Normal"]
    story = []
    for line in text.split("\n"):
        safe = (line
                .replace("&", "&amp;")
                .replace("<", "&lt;")
                .replace(">", "&gt;")) or "&nbsp;"
        story.append(Paragraph(safe, style))
        story.append(Spacer(1, 3))
    doc.build(story)
    buf.seek(0)
    return buf.read(), "converted.pdf", MIME["PDF"]


def txt_to_docx(data):
    from docx import Document
    text = data.decode("utf-8", errors="replace")
    doc = Document()
    for line in text.split("\n"):
        doc.add_paragraph(line)
    out = io.BytesIO()
    doc.save(out)
    out.seek(0)
    return out.read(), "converted.docx", MIME["DOCX"]


# ─── Spreadsheet / data conversions ────────────────────────────────────────────
def csv_to_xlsx(data):
    import pandas as pd
    df = pd.read_csv(io.BytesIO(data))
    buf = io.BytesIO()
    df.to_excel(buf, index=False, engine="openpyxl")
    buf.seek(0)
    return buf.read(), "converted.xlsx", MIME["XLSX"]


def xlsx_to_csv(data):
    import pandas as pd
    df = pd.read_excel(io.BytesIO(data), engine="openpyxl")
    buf = io.BytesIO()
    df.to_csv(buf, index=False)
    buf.seek(0)
    return buf.read(), "converted.csv", MIME["CSV"]


def csv_to_json(data):
    import pandas as pd
    df = pd.read_csv(io.BytesIO(data))
    return (
        df.to_json(orient="records", indent=2).encode("utf-8"),
        "converted.json",
        MIME["JSON"],
    )


def xlsx_to_json(data):
    import pandas as pd
    df = pd.read_excel(io.BytesIO(data), engine="openpyxl")
    return (
        df.to_json(orient="records", indent=2).encode("utf-8"),
        "converted.json",
        MIME["JSON"],
    )


def json_to_csv(data):
    import pandas as pd
    df = pd.read_json(io.BytesIO(data))
    buf = io.BytesIO()
    df.to_csv(buf, index=False)
    buf.seek(0)
    return buf.read(), "converted.csv", MIME["CSV"]


def json_to_xlsx(data):
    import pandas as pd
    df = pd.read_json(io.BytesIO(data))
    buf = io.BytesIO()
    df.to_excel(buf, index=False, engine="openpyxl")
    buf.seek(0)
    return buf.read(), "converted.xlsx", MIME["XLSX"]


# ─── Router ─────────────────────────────────────────────────────────────────────
IMAGES = {"PNG", "JPG", "WEBP", "BMP"}


def convert(data, src, dst):
    if src == "PDF":
        if dst == "DOCX": return pdf_to_docx(data)
        if dst == "TXT":  return pdf_to_txt(data)
        if dst == "PPTX": return pdf_to_pptx(data)
        if dst in IMAGES: return pdf_to_image(data, dst)

    elif src == "DOCX":
        if dst == "PDF": return docx_to_pdf(data)
        if dst == "TXT": return docx_to_txt(data)

    elif src == "PPTX":
        if dst == "PDF":    return pptx_to_pdf(data)
        if dst in IMAGES:   return pptx_to_image(data, dst)

    elif src in IMAGES:
        if dst == "PDF":    return img_to_pdf(data)
        if dst == "DOCX":   return img_to_docx(data)
        if dst in IMAGES:   return img_to_img(data, dst)

    elif src == "TXT":
        if dst == "PDF":  return txt_to_pdf(data)
        if dst == "DOCX": return txt_to_docx(data)

    elif src == "CSV":
        if dst == "XLSX": return csv_to_xlsx(data)
        if dst == "JSON": return csv_to_json(data)

    elif src == "XLSX":
        if dst == "CSV":  return xlsx_to_csv(data)
        if dst == "JSON": return xlsx_to_json(data)

    elif src == "JSON":
        if dst == "CSV":  return json_to_csv(data)
        if dst == "XLSX": return json_to_xlsx(data)

    raise ValueError(f"No handler for {src} → {dst}")


# ─── Dependency checker ─────────────────────────────────────────────────────────
def check_deps(src, dst):
    """Returns (ok: bool, missing: list[str])"""
    needed = REQUIRES.get((src, dst), [])
    missing = []
    for dep in needed:
        if dep == "LibreOffice":
            if not _find_libreoffice():
                missing.append("LibreOffice (system install)")
        elif dep == "PyMuPDF":
            if not HAS_PYMUPDF: missing.append("PyMuPDF  (`pip install PyMuPDF`)")
        elif dep == "pdf2docx":
            if not HAS_PDF2DOCX: missing.append("pdf2docx  (`pip install pdf2docx`)")
        elif dep == "python-pptx":
            if not HAS_PPTX: missing.append("python-pptx  (`pip install python-pptx`)")
        elif dep == "python-docx":
            if not HAS_DOCX: missing.append("python-docx  (`pip install python-docx`)")
        elif dep == "Pillow":
            if not HAS_PIL: missing.append("Pillow  (`pip install Pillow`)")
        elif dep == "reportlab":
            if not HAS_REPORTLAB: missing.append("reportlab  (`pip install reportlab`)")
        elif dep == "pandas":
            if not HAS_PANDAS: missing.append("pandas  (`pip install pandas`)")
        elif dep == "openpyxl":
            if not HAS_PANDAS: missing.append("openpyxl  (`pip install openpyxl`)")
    return len(missing) == 0, missing


# ═══════════════════════════════════════════════════════════════════════════════
# Streamlit UI
# ═══════════════════════════════════════════════════════════════════════════════
st.set_page_config(
    page_title="File Converter",
    page_icon="⚡",
    layout="centered",
    initial_sidebar_state="collapsed",
)

# ── Custom CSS ──────────────────────────────────────────────────────────────────
st.markdown("""
<style>
    @import url('https://fonts.googleapis.com/css2?family=DM+Mono:wght@400;500&family=DM+Sans:wght@300;400;500;600&display=swap');

    html, body, [class*="css"] { font-family: 'DM Sans', sans-serif; }

    .block-container {
        max-width: 660px;
        padding-top: 2.5rem;
        padding-bottom: 3rem;
    }

    /* Header */
    .app-header {
        text-align: center;
        margin-bottom: 2.5rem;
    }
    .app-header h1 {
        font-size: 2rem;
        font-weight: 600;
        letter-spacing: -0.03em;
        margin: 0;
        color: #0f0f0f;
    }
    .app-header p {
        color: #6b7280;
        font-size: 0.9rem;
        margin-top: 0.4rem;
    }

    /* Format cards */
    .fmt-card {
        background: #f9fafb;
        border: 1.5px solid #e5e7eb;
        border-radius: 12px;
        padding: 1.2rem 1.4rem;
        transition: border-color 0.15s;
    }

    /* Arrow indicator */
    .arrow-col {
        display: flex;
        align-items: flex-end;
        justify-content: center;
        padding-bottom: 0.6rem;
        font-size: 1.6rem;
        color: #9ca3af;
    }

    /* Step labels */
    .step-label {
        font-size: 0.7rem;
        font-weight: 600;
        letter-spacing: 0.08em;
        text-transform: uppercase;
        color: #9ca3af;
        margin-bottom: 0.3rem;
    }

    /* Upload zone */
    .upload-zone {
        border: 2px dashed #d1d5db;
        border-radius: 12px;
        padding: 1rem;
        margin-top: 1rem;
        background: #fafafa;
    }

    /* Selectbox styling */
    div[data-testid="stSelectbox"] > label {
        font-size: 0.72rem !important;
        font-weight: 600 !important;
        text-transform: uppercase;
        letter-spacing: 0.06em;
        color: #6b7280 !important;
    }
    div[data-testid="stSelectbox"] > div > div {
        border-radius: 10px !important;
        border: 1.5px solid #e5e7eb !important;
        font-size: 1rem !important;
        font-weight: 500 !important;
    }

    /* Buttons */
    .stButton > button, .stDownloadButton > button {
        border-radius: 10px !important;
        font-weight: 500 !important;
        letter-spacing: 0.01em !important;
        transition: all 0.15s !important;
    }
    .stButton > button[kind="primary"] {
        background: #111827 !important;
        border: none !important;
    }
    .stButton > button[kind="primary"]:hover {
        background: #374151 !important;
        transform: translateY(-1px);
        box-shadow: 0 4px 12px rgba(0,0,0,0.15) !important;
    }
    .stDownloadButton > button {
        background: #059669 !important;
        color: white !important;
        border: none !important;
    }
    .stDownloadButton > button:hover {
        background: #047857 !important;
        transform: translateY(-1px);
        box-shadow: 0 4px 12px rgba(5,150,105,0.3) !important;
    }

    /* File uploader */
    div[data-testid="stFileUploader"] > label {
        font-size: 0.72rem !important;
        font-weight: 600 !important;
        text-transform: uppercase;
        letter-spacing: 0.06em;
        color: #6b7280 !important;
    }
    div[data-testid="stFileUploadDropzone"] {
        border-radius: 12px !important;
        border: 2px dashed #d1d5db !important;
    }

    /* Divider */
    hr { margin: 1.5rem 0; border-color: #f3f4f6; }

    /* Alerts */
    div[data-testid="stAlert"] { border-radius: 10px !important; }

    /* Expander */
    div[data-testid="stExpander"] summary {
        font-size: 0.85rem;
        font-weight: 500;
        color: #6b7280;
    }

    /* Footer */
    footer { visibility: hidden; }
    #MainMenu { visibility: hidden; }
    header { visibility: hidden; }
</style>
""", unsafe_allow_html=True)

# ── Header ──────────────────────────────────────────────────────────────────────
st.markdown("""
<div class="app-header">
    <h1>⚡ File Converter</h1>
    <p>Convert between 11 formats — documents, images, spreadsheets</p>
</div>
""", unsafe_allow_html=True)

# ── Step 1 & 2: Format selection ────────────────────────────────────────────────
col1, arrow, col2 = st.columns([5, 1, 5])

with col1:
    from_fmt = st.selectbox(
        "I HAVE",
        list(CONVERSIONS.keys()),
        format_func=lambda f: f"{ICONS[f]}  {f}",
        key=f"from_fmt",
    )

with arrow:
    st.markdown("""
    <div class="arrow-col">→</div>
    """, unsafe_allow_html=True)

with col2:
    to_opts = CONVERSIONS[from_fmt]
    to_fmt = st.selectbox(
        "I WANT",
        to_opts,
        format_func=lambda f: f"{ICONS[f]}  {f}",
        key=f"to_fmt_{from_fmt}",  # key changes with from_fmt → auto-resets
    )

st.write("")

# ── Dependency check ─────────────────────────────────────────────────────────────
deps_ok, missing_deps = check_deps(from_fmt, to_fmt)
if not deps_ok:
    st.warning(
        "⚠️  Missing dependencies for this conversion:\n\n" +
        "\n".join(f"• {d}" for d in missing_deps)
    )

# ── Step 3: File upload ──────────────────────────────────────────────────────────
uploaded = st.file_uploader(
    "UPLOAD FILE",
    type=ACCEPT[from_fmt],
    key=f"upload_{from_fmt}",  # resets when format changes
    help=f"Accepted: {', '.join('.' + e for e in ACCEPT[from_fmt])}",
    label_visibility="visible",
)

if uploaded:
    data = uploaded.read()
    size_kb = len(data) / 1024
    size_str = f"{size_kb:.1f} KB" if size_kb < 1024 else f"{size_kb/1024:.2f} MB"
    st.caption(f"📎  **{uploaded.name}**  ·  {size_str}")
    st.write("")

    btn_label = f"Convert  {ICONS[from_fmt]} {from_fmt}  →  {ICONS[to_fmt]} {to_fmt}"
    if st.button(btn_label, type="primary", use_container_width=True, disabled=not deps_ok):
        with st.spinner(f"Converting {from_fmt} → {to_fmt}…"):
            try:
                out_data, out_name, out_mime = convert(data, from_fmt, to_fmt)
                out_kb = len(out_data) / 1024
                out_size_str = f"{out_kb:.1f} KB" if out_kb < 1024 else f"{out_kb/1024:.2f} MB"

                st.success(f"✅  Done!  Output: **{out_name}**  ·  {out_size_str}")

                # Special note for multi-page ZIP output
                if out_name.endswith(".zip"):
                    st.info(
                        "ℹ️  Multi-page document detected — pages are packaged as a ZIP file. "
                        "Extract it to access individual images."
                    )

                st.download_button(
                    label=f"⬇️  Download {out_name}",
                    data=out_data,
                    file_name=out_name,
                    mime=out_mime,
                    use_container_width=True,
                )
            except Exception as e:
                st.error(f"❌  Conversion failed:\n\n{e}")

# ── Divider + Info sections ─────────────────────────────────────────────────────
st.divider()

col_a, col_b = st.columns(2)

with col_a:
    with st.expander("📋  All supported conversions"):
        for src, dsts in CONVERSIONS.items():
            targets = "  ·  ".join(f"{ICONS[d]} {d}" for d in dsts)
            st.markdown(f"**{ICONS[src]} {src}** → {targets}")

with col_b:
    with st.expander("📦  Setup & requirements"):
        st.markdown("**Python packages:**")
        st.code(
            "pip install streamlit PyMuPDF pdf2docx \\\n"
            "    python-pptx python-docx Pillow \\\n"
            "    reportlab pandas openpyxl",
            language="bash",
        )
        st.markdown("**For DOCX/PPTX → PDF:** LibreOffice required")
        st.code(
            "# macOS\nbrew install --cask libreoffice\n\n"
            "# Ubuntu/Debian\nsudo apt install libreoffice\n\n"
            "# Windows: libreoffice.org/download",
            language="bash",
        )
